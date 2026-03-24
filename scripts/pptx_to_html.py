"""PPTX to HTML converter - preserves layout, images, tables, and text styling."""

import sys
import base64
import html
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def emu_to_px(emu):
    """Convert EMU to pixels (96 DPI)."""
    return round(emu / 914400 * 96)


def color_to_css(color_obj):
    """Extract CSS color string from a pptx color object."""
    try:
        if color_obj and color_obj.type is not None:
            rgb = color_obj.rgb
            if rgb:
                return f"#{rgb}"
    except Exception:
        pass
    return None


def fill_to_css(fill):
    """Extract background color from shape fill."""
    try:
        if fill and fill.type is not None:
            fg = fill.fore_color
            if fg and fg.type is not None:
                rgb = fg.rgb
                if rgb:
                    return f"#{rgb}"
    except Exception:
        pass
    return None


def get_paragraph_html(paragraph):
    """Convert a paragraph to HTML with inline styling."""
    runs_html = []
    for run in paragraph.runs:
        text = html.escape(run.text)
        if not text:
            continue

        styles = []
        font = run.font

        # Font size
        if font.size:
            styles.append(f"font-size:{emu_to_px(font.size)}px")

        # Bold
        if font.bold:
            styles.append("font-weight:bold")

        # Italic
        if font.italic:
            styles.append("font-style:italic")

        # Underline
        if font.underline:
            styles.append("text-decoration:underline")

        # Font color
        color = color_to_css(font.color)
        if color:
            styles.append(f"color:{color}")

        # Font family
        if font.name:
            styles.append(f"font-family:'{font.name}',sans-serif")

        if styles:
            runs_html.append(f'<span style="{";".join(styles)}">{text}</span>')
        else:
            runs_html.append(text)

    return "".join(runs_html)


def get_alignment_css(paragraph):
    """Get CSS text-align from paragraph alignment."""
    align_map = {
        PP_ALIGN.LEFT: "left",
        PP_ALIGN.CENTER: "center",
        PP_ALIGN.RIGHT: "right",
        PP_ALIGN.JUSTIFY: "justify",
    }
    if paragraph.alignment and paragraph.alignment in align_map:
        return align_map[paragraph.alignment]
    return None


def render_text_frame(text_frame):
    """Render a text frame to HTML."""
    paras = []
    for para in text_frame.paragraphs:
        p_html = get_paragraph_html(para)
        if not p_html.strip():
            paras.append("<p>&nbsp;</p>")
            continue
        align = get_alignment_css(para)
        style = f' style="text-align:{align}"' if align else ""
        paras.append(f"<p{style}>{p_html}</p>")
    return "\n".join(paras)


def render_table(table):
    """Render a table to HTML."""
    rows_html = []
    for row in table.rows:
        cells_html = []
        for cell in row.cells:
            # Cell background
            cell_styles = []
            bg = fill_to_css(cell.fill)
            if bg:
                cell_styles.append(f"background:{bg}")

            content = render_text_frame(cell.text_frame)
            style_attr = f' style="{";".join(cell_styles)}"' if cell_styles else ""

            # Check for merged cells
            span_attrs = ""
            if cell.span_height > 1:
                span_attrs += f' rowspan="{cell.span_height}"'
            if cell.span_width > 1:
                span_attrs += f' colspan="{cell.span_width}"'

            # Skip cells that are merged into another
            if cell.is_merge_origin or (cell.span_height == 1 and cell.span_width == 1):
                cells_html.append(f"<td{span_attrs}{style_attr}>{content}</td>")

        rows_html.append(f"<tr>{''.join(cells_html)}</tr>")
    return f'<table class="pptx-table">{"".join(rows_html)}</table>'


def render_image(shape):
    """Render an image as base64 embedded img tag."""
    try:
        image = shape.image
        content_type = image.content_type
        blob = image.blob
        b64 = base64.b64encode(blob).decode("utf-8")
        return f'<img src="data:{content_type};base64,{b64}" />'
    except Exception:
        return ""


def render_shape(shape):
    """Render a single shape to positioned HTML div."""
    left = emu_to_px(shape.left)
    top = emu_to_px(shape.top)
    width = emu_to_px(shape.width)
    height = emu_to_px(shape.height)

    # Shape background & border
    extra_styles = []
    try:
        bg = fill_to_css(shape.fill)
        if bg:
            extra_styles.append(f"background:{bg}")
    except Exception:
        pass

    try:
        line = shape.line
        if line and line.fill and line.fill.type is not None:
            line_color = color_to_css(line.color)
            line_width = emu_to_px(line.width) if line.width else 1
            if line_color:
                extra_styles.append(f"border:{line_width}px solid {line_color}")
    except Exception:
        pass

    # Rotation
    rotation = ""
    try:
        if shape.rotation and shape.rotation != 0:
            rotation = f"transform:rotate({shape.rotation}deg);"
    except Exception:
        pass

    # Content
    content = ""
    shape_type = shape.shape_type

    # Image
    if shape_type == 13:  # PICTURE
        content = render_image(shape)

    # Table
    elif shape_type == 19:  # TABLE
        content = render_table(shape.table)

    # Text-bearing shapes
    elif hasattr(shape, "has_text_frame") and shape.has_text_frame:
        content = render_text_frame(shape.text_frame)

    if not content.strip() and not extra_styles:
        # Skip empty shapes with no visual styling
        return ""

    style = (
        f"position:absolute;left:{left}px;top:{top}px;"
        f"width:{width}px;height:{height}px;"
        f"overflow:hidden;{rotation}"
    )
    if extra_styles:
        style += ";".join(extra_styles) + ";"

    return f'<div class="shape" style="{style}">{content}</div>'


def convert_pptx_to_html(pptx_path, output_path=None):
    """Main conversion function."""
    pptx_path = Path(pptx_path)
    if output_path is None:
        output_path = pptx_path.with_suffix(".html")
    else:
        output_path = Path(output_path)

    prs = Presentation(str(pptx_path))
    slide_w = emu_to_px(prs.slide_width)
    slide_h = emu_to_px(prs.slide_height)

    slides_html = []
    for i, slide in enumerate(prs.slides):
        # Slide background
        slide_bg = ""
        try:
            bg_fill = slide.background.fill
            bg_color = fill_to_css(bg_fill)
            if bg_color:
                slide_bg = f"background:{bg_color};"
        except Exception:
            pass

        shapes_html = []
        for shape in slide.shapes:
            rendered = render_shape(shape)
            if rendered:
                shapes_html.append(rendered)

        slides_html.append(
            f'<div class="slide" style="width:{slide_w}px;height:{slide_h}px;{slide_bg}">'
            f'<div class="slide-label">Slide {i+1}</div>'
            f'{"".join(shapes_html)}'
            f"</div>"
        )

    title = html.escape(pptx_path.stem)

    full_html = f"""<!DOCTYPE html>
<html lang="ko">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{title}</title>
<style>
* {{ margin:0; padding:0; box-sizing:border-box; }}
body {{
    background:#f0f0f0;
    font-family:'Malgun Gothic','맑은 고딕',sans-serif;
    display:flex;
    flex-direction:column;
    align-items:center;
    padding:20px;
    gap:30px;
}}
.slide {{
    position:relative;
    background:#fff;
    box-shadow:0 4px 20px rgba(0,0,0,0.15);
    overflow:hidden;
    page-break-after:always;
}}
.slide-label {{
    position:absolute;
    top:8px;
    right:12px;
    background:rgba(0,0,0,0.5);
    color:#fff;
    padding:2px 10px;
    border-radius:10px;
    font-size:12px;
    z-index:1000;
}}
.shape p {{
    margin:0;
    padding:1px 0;
    line-height:1.3;
}}
.shape img {{
    width:100%;
    height:100%;
    object-fit:contain;
}}
.pptx-table {{
    width:100%;
    height:100%;
    border-collapse:collapse;
    font-size:11px;
}}
.pptx-table td {{
    border:1px solid #999;
    padding:3px 5px;
    vertical-align:middle;
}}
.pptx-table p {{
    margin:0;
    line-height:1.2;
}}
@media print {{
    body {{ background:white; padding:0; gap:0; }}
    .slide {{ box-shadow:none; }}
    .slide-label {{ display:none; }}
}}
</style>
</head>
<body>
{"".join(slides_html)}
</body>
</html>"""

    output_path.write_text(full_html, encoding="utf-8")
    print(f"Converted: {output_path}")
    return str(output_path)


if __name__ == "__main__":
    src = sys.argv[1] if len(sys.argv) > 1 else None
    dst = sys.argv[2] if len(sys.argv) > 2 else None
    if not src:
        print("Usage: python pptx_to_html.py <input.pptx> [output.html]")
        sys.exit(1)
    convert_pptx_to_html(src, dst)
